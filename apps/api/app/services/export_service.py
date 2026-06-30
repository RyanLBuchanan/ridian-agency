"""Artifact export + open helpers.

Three responsibilities:

1. **Path validation** — every operation resolves the operator-supplied
   folder and refuses anything that doesn't sit inside the configured
   ``outputs/`` directory. Filenames are checked against a small allowlist.

2. **Open** — ``os.startfile`` opens a folder in Explorer or a file in its
   registered default app (Word for ``.docx``, PowerPoint for ``.pptx``,
   text editor for ``.md``, etc.). Windows-only for now; other platforms
   return a clear "not supported" error.

3. **Export** — ZIP, DOCX (from ``business_document.md``), and PPTX
   (from ``slide_outline.md``). Each export is written next to the
   artifact folder (ZIP) or inside it (DOCX/PPTX).
"""

from __future__ import annotations

import logging
import os
import re
import shutil
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

from docx import Document
from pptx import Presentation

from .artifact_service import outputs_dir

log = logging.getLogger("ridian.export")

# Files the operator is allowed to open via the API. Anything else is
# rejected even if it lives inside the artifact folder.
ALLOWED_OPEN_FILENAMES: frozenset[str] = frozenset(
    {
        "task.txt",
        # Business workflow artifacts
        "research_summary.md",
        "business_document.md",
        "slide_outline.md",
        "draft_email.md",
        # Files we generate ourselves and want to be openable.
        "business_document.docx",
        "slide_outline.pptx",
        # Social Media Production workflow artifacts
        "social_content_package.md",
        "script.md",
        "caption_package.md",
        "posting_checklist.md",
        "visual_production.md",
        # Agentic Advances Daily Brief
        "agentic_advances_brief.md",
        # NotebookLM Package
        "notebooklm_package.md",
        # Operator v1 — audiobook + sources + script + machine log
        "sources_packet.md",
        "script.md",
        "audiobook.mp3",
        "operation_log.json",
        # NotebookLM research packet (build_research_packet)
        "research_packet.md",
        # General prose deliverable written via write_file (title is the H1).
        "document.md",
    }
)


@dataclass
class ExportError(Exception):
    """Raised by export helpers when a request is invalid or fails.

    ``status`` maps to the HTTP status the API layer should return.
    Detail is safe to surface to the operator (never includes secrets).
    """

    detail: str
    status: int = 400

    def __str__(self) -> str:  # for default Exception formatting
        return self.detail


# ---------------------------------------------------------------------------
# Path validation
# ---------------------------------------------------------------------------


def _resolve_artifact_folder(folder_str: str) -> Path:
    """Resolve and validate that ``folder_str`` is a real folder inside outputs."""
    if not folder_str or not isinstance(folder_str, str):
        raise ExportError("artifact_folder is required.", status=400)

    try:
        candidate = Path(folder_str).resolve(strict=False)
    except (OSError, ValueError) as exc:
        raise ExportError(f"Invalid folder path ({type(exc).__name__}).", status=400) from exc

    outputs = outputs_dir().resolve()

    # is_relative_to (3.9+) rejects path traversal and any folder outside outputs.
    try:
        candidate.relative_to(outputs)
    except ValueError as exc:
        raise ExportError(
            "Folder is not inside the configured outputs directory.",
            status=400,
        ) from exc

    if not candidate.exists():
        raise ExportError("Artifact folder does not exist.", status=404)
    if not candidate.is_dir():
        raise ExportError("Artifact path is not a directory.", status=400)
    # Reject the outputs root itself — must be a per-run subfolder.
    if candidate.resolve() == outputs:
        raise ExportError(
            "Refusing to operate on the outputs root. Provide a per-run folder.",
            status=400,
        )

    return candidate


def _resolve_artifact_file(folder_str: str, filename: str) -> Path:
    """Validate folder + filename, return the resolved file path."""
    if not filename or not isinstance(filename, str):
        raise ExportError("filename is required.", status=400)
    if filename not in ALLOWED_OPEN_FILENAMES:
        raise ExportError(
            f"Filename {filename!r} is not in the allowlist.",
            status=400,
        )

    folder = _resolve_artifact_folder(folder_str)
    candidate = (folder / filename).resolve(strict=False)

    # Defense in depth: the resolved file must still be inside the folder.
    try:
        candidate.relative_to(folder)
    except ValueError as exc:
        raise ExportError("Resolved file escapes the artifact folder.", status=400) from exc

    if not candidate.exists():
        raise ExportError(f"File {filename!r} does not exist in this run.", status=404)
    if not candidate.is_file():
        raise ExportError(f"{filename!r} is not a regular file.", status=400)

    return candidate


# ---------------------------------------------------------------------------
# Open in default app
# ---------------------------------------------------------------------------


def _native_open(path: Path) -> None:
    """Open ``path`` in the OS default handler. Windows uses os.startfile.

    On macOS / Linux we'd shell out to `open` / `xdg-open`. Not implemented
    here — the desktop app is Windows-only for v1, and we'd rather refuse
    explicitly than silently no-op.
    """
    if sys.platform.startswith("win"):
        os.startfile(str(path))  # type: ignore[attr-defined]  # win-only
        return
    raise ExportError(
        f"Opening files from the API is only implemented on Windows (got {sys.platform}).",
        status=501,
    )


def open_artifact_folder(folder_str: str) -> Path:
    folder = _resolve_artifact_folder(folder_str)
    _native_open(folder)
    log.info("artifact.open_folder folder=%s", folder)
    return folder


def open_artifact_file(folder_str: str, filename: str) -> Path:
    path = _resolve_artifact_file(folder_str, filename)
    _native_open(path)
    log.info("artifact.open_file file=%s", path)
    return path


# ---------------------------------------------------------------------------
# ZIP export
# ---------------------------------------------------------------------------


def export_zip(folder_str: str) -> Path:
    """Zip an entire artifact folder. ZIP lands next to the folder under outputs/."""
    folder = _resolve_artifact_folder(folder_str)
    outputs = outputs_dir().resolve()

    base_name = folder.name
    zip_base = outputs / base_name  # shutil appends '.zip'

    # shutil.make_archive overwrites if the target exists.
    archive = shutil.make_archive(
        base_name=str(zip_base),
        format="zip",
        root_dir=str(folder.parent),
        base_dir=folder.name,
    )
    zip_path = Path(archive).resolve()
    log.info("artifact.export_zip folder=%s zip=%s", folder, zip_path)
    return zip_path


# ---------------------------------------------------------------------------
# Markdown -> DOCX
# ---------------------------------------------------------------------------

_HEADING_RE = re.compile(r"^(#{1,4})\s+(.*)$")
_BULLET_RE = re.compile(r"^\s*[-*]\s+(.*)$")
_NUMBERED_RE = re.compile(r"^\s*\d+\.\s+(.*)$")
_HR_RE = re.compile(r"^\s*---+\s*$")


def _strip_inline_markdown(s: str) -> str:
    """Drop **bold**, *italic*, and `code` markers — keep plain text."""
    s = re.sub(r"\*\*([^*]+)\*\*", r"\1", s)
    s = re.sub(r"(?<!\*)\*([^*\n]+)\*(?!\*)", r"\1", s)
    s = re.sub(r"`([^`]+)`", r"\1", s)
    return s


def _build_docx_from_markdown(md_text: str) -> Document:
    """Convert a Markdown string to a python-docx Document.

    Supports headings (#-####), bullet lists (- / *), numbered lists,
    paragraphs, and horizontal rules. Inline emphasis is flattened to
    plain text — readable, predictable, no surprise styling.
    """
    doc = Document()
    lines = md_text.replace("\r\n", "\n").split("\n")
    i = 0

    def add_paragraph(text: str) -> None:
        if text:
            doc.add_paragraph(_strip_inline_markdown(text))

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if not stripped:
            i += 1
            continue

        if _HR_RE.match(line):
            # No native HR in python-docx; insert a blank paragraph as a separator.
            doc.add_paragraph("")
            i += 1
            continue

        m = _HEADING_RE.match(stripped)
        if m:
            level = len(m.group(1))
            doc.add_heading(_strip_inline_markdown(m.group(2).strip()), level=level)
            i += 1
            continue

        if _BULLET_RE.match(line):
            while i < len(lines) and _BULLET_RE.match(lines[i]):
                item = _BULLET_RE.match(lines[i]).group(1)
                doc.add_paragraph(_strip_inline_markdown(item), style="List Bullet")
                i += 1
            continue

        if _NUMBERED_RE.match(line):
            while i < len(lines) and _NUMBERED_RE.match(lines[i]):
                item = _NUMBERED_RE.match(lines[i]).group(1)
                doc.add_paragraph(_strip_inline_markdown(item), style="List Number")
                i += 1
            continue

        # Plain paragraph — collect consecutive non-blank, non-special lines.
        buf: list[str] = []
        while (
            i < len(lines)
            and lines[i].strip()
            and not _HEADING_RE.match(lines[i].strip())
            and not _BULLET_RE.match(lines[i])
            and not _NUMBERED_RE.match(lines[i])
            and not _HR_RE.match(lines[i])
        ):
            buf.append(lines[i].strip())
            i += 1
        add_paragraph(" ".join(buf))

    return doc


def export_docx(folder_str: str) -> Path:
    """Build a .docx from business_document.md inside the given run folder."""
    folder = _resolve_artifact_folder(folder_str)
    source = folder / "business_document.md"
    if not source.exists():
        raise ExportError("business_document.md not found in this run.", status=404)

    md_text = source.read_text(encoding="utf-8")
    doc = _build_docx_from_markdown(md_text)
    out_path = folder / "business_document.docx"
    doc.save(str(out_path))
    log.info("artifact.export_docx out=%s", out_path)
    return out_path


# ---------------------------------------------------------------------------
# Markdown slide outline -> PPTX
# ---------------------------------------------------------------------------

# Matches "## Slide 1 — Title" with either em-dash, en-dash, hyphen, or colon.
_SLIDE_HEADING_RE = re.compile(
    r"^##\s*Slide\s*(\d+)\s*[—–\-:]\s*(.+?)\s*$",
    re.IGNORECASE,
)
_SPEAKER_NOTE_RE = re.compile(r"^Speaker\s*note\s*:\s*(.*)$", re.IGNORECASE)


def _parse_slide_blocks(md_text: str) -> list[dict[str, object]]:
    """Split the slide outline into per-slide blocks."""
    blocks: list[dict[str, object]] = []
    current: Optional[dict[str, object]] = None

    for raw in md_text.replace("\r\n", "\n").split("\n"):
        line = raw.rstrip()

        m = _SLIDE_HEADING_RE.match(line)
        if m:
            if current is not None:
                blocks.append(current)
            current = {"title": _strip_inline_markdown(m.group(2)), "bullets": [], "notes": ""}
            continue

        if current is None:
            continue

        if _HR_RE.match(line):
            continue

        if _BULLET_RE.match(line):
            item = _BULLET_RE.match(line).group(1)
            current["bullets"].append(_strip_inline_markdown(item))  # type: ignore[index]
            continue

        note_match = _SPEAKER_NOTE_RE.match(line.strip())
        if note_match:
            note_text = _strip_inline_markdown(note_match.group(1).strip())
            existing = current.get("notes", "")
            current["notes"] = (existing + "\n" + note_text).strip() if existing else note_text
            continue

        # Plain text under a slide gets appended to notes too — better than dropping it.
        if line.strip():
            existing = current.get("notes", "")
            extra = _strip_inline_markdown(line.strip())
            current["notes"] = (existing + "\n" + extra).strip() if existing else extra

    if current is not None:
        blocks.append(current)
    return blocks


def _build_pptx_from_outline(md_text: str) -> Presentation:
    prs = Presentation()
    # Layout 1 is "Title and Content" in the default template.
    layout = prs.slide_layouts[1]

    blocks = _parse_slide_blocks(md_text)
    if not blocks:
        # Emit a single placeholder slide so the user gets a valid file even
        # if the outline didn't match our expected format.
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = "Slide outline"
        body = slide.placeholders[1].text_frame
        body.text = "No slide blocks were recognized in slide_outline.md."
        return prs

    for block in blocks:
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = str(block.get("title") or "Slide")

        body_tf = slide.placeholders[1].text_frame
        bullets = block.get("bullets") or []  # type: ignore[assignment]
        if bullets:
            body_tf.text = bullets[0]
            for b in bullets[1:]:
                p = body_tf.add_paragraph()
                p.text = b
        else:
            body_tf.text = ""

        notes_text = str(block.get("notes") or "")
        if notes_text:
            slide.notes_slide.notes_text_frame.text = notes_text

    return prs


def export_pptx(folder_str: str) -> Path:
    """Build a .pptx from slide_outline.md inside the given run folder."""
    folder = _resolve_artifact_folder(folder_str)
    source = folder / "slide_outline.md"
    if not source.exists():
        raise ExportError("slide_outline.md not found in this run.", status=404)

    md_text = source.read_text(encoding="utf-8")
    prs = _build_pptx_from_outline(md_text)
    out_path = folder / "slide_outline.pptx"
    prs.save(str(out_path))
    log.info("artifact.export_pptx out=%s", out_path)
    return out_path


# TODO(google-workspace): add upload_to_drive(), export_google_doc(),
# export_google_slides() once OAuth wiring lands. See the
# "Future Google Workspace exports" section in README.md for the design.
# Implementation notes:
#   - OAuth tokens go in apps/api/local_settings.json (same security model
#     as smtp_password and openai_api_key — never returned by any endpoint).
#   - Use the narrow drive.file scope, never full drive.
#   - Each function should mirror the export_zip/docx/pptx signature:
#     take artifact_folder, return the remote URL/id.
